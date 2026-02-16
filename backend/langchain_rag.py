import os
from typing import List, Optional
from langchain_groq import ChatGroq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.documents import Document
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from dotenv import load_dotenv
from pypdf import PdfReader
from pptx import Presentation

load_dotenv()


def sync_documents_from_blob(local_folder: str = "documents"):
    """Download documents from Azure Blob Storage to local folder.
    
    Requires env vars:
    - AZURE_STORAGE_CONNECTION_STRING: From Azure Portal → Storage Account → Access keys
    - AZURE_STORAGE_CONTAINER_NAME: Blob container name (default: 'documents')
    
    Falls back to local documents/ folder if not configured (for local dev).
    """
    connection_string = os.environ.get("AZURE_STORAGE_CONNECTION_STRING")
    container_name = os.environ.get("AZURE_STORAGE_CONTAINER_NAME", "documents")
    
    if not connection_string:
        print("ℹ️  AZURE_STORAGE_CONNECTION_STRING not set — using local documents/ folder")
        return
    
    try:
        from azure.storage.blob import ContainerClient
        
        os.makedirs(local_folder, exist_ok=True)
        container_client = ContainerClient.from_connection_string(connection_string, container_name)
        
        blob_list = list(container_client.list_blobs())
        blob_names = {blob.name for blob in blob_list}
        print(f"📥 Found {len(blob_list)} file(s) in Azure Blob Storage '{container_name}'")
        
        # Delete local files that no longer exist in blob storage
        if os.path.exists(local_folder):
            for local_file in os.listdir(local_folder):
                if local_file not in blob_names:
                    local_path = os.path.join(local_folder, local_file)
                    print(f"  🗑️  Deleting orphaned file: {local_file}")
                    os.remove(local_path)
        
        for blob in blob_list:
            local_path = os.path.join(local_folder, blob.name)
            # Skip if already downloaded and same size
            if os.path.exists(local_path) and os.path.getsize(local_path) == blob.size:
                print(f"  ✓ Already up to date: {blob.name}")
                continue
            
            print(f"  ⬇ Downloading: {blob.name} ({blob.size} bytes)")
            blob_data = container_client.download_blob(blob.name).readall()
            with open(local_path, "wb") as f:
                f.write(blob_data)
        
        print("✓ Azure Blob Storage sync complete")
    except Exception as e:
        print(f"⚠️ Azure Blob sync failed: {e}. Falling back to local documents/")


class LangChainRAG:
    # LangChain-based RAG using Groq (LLM) and HuggingFace (Embeddings)
    def __init__(self, documents_folder: str = "documents"):
        self.documents_folder, self.groq_api_key = documents_folder, os.environ.get("GROQ_API_KEY")
        if not self.groq_api_key: raise ValueError("GROQ_API_KEY missing!")
        
        # Sync documents from Azure Blob Storage (skipped if not configured)
        sync_documents_from_blob(self.documents_folder)
        
        # Initialize embeddings (HuggingFace, local, free)
        from langchain_community.embeddings import HuggingFaceEmbeddings
        self.embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2", model_kwargs={'device': 'cpu'})
        
        # Initialize LLM (Groq Llama 3.3, fast, free tier)
        self.llm = ChatGroq(model="llama-3.3-70b-versatile", api_key=self.groq_api_key, temperature=0.7)
        print("✓ Groq API initialized")
        
        self.vectorstore, self.rag_chain = None, None
        
    def load_documents(self) -> List[Document]:
        # Re-sync from blob storage before loading
        sync_documents_from_blob(self.documents_folder)
        
        documents = []
        if not os.path.exists(self.documents_folder):
            os.makedirs(self.documents_folder)
            print(f"Created {self.documents_folder} folder.")
            return documents
        
        for filename in os.listdir(self.documents_folder):
            filepath = os.path.join(self.documents_folder, filename)
            try:
                text = ""
                # Load PDF
                if filename.endswith(".pdf"):
                    for page in PdfReader(filepath).pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    if text.strip():
                        documents.append(Document(page_content=text.strip(), metadata={"source": filename, "type": "pdf"}))
                        print(f"Loaded: {filename} ({len(text)} chars)")
                    else:
                        print(f"⚠️ Skipping {filename}: no extractable text")
                # Load PPTX
                elif filename.endswith(".pptx"):
                    for slide in Presentation(filepath).slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                text += shape.text + "\n"
                    if text.strip():
                        documents.append(Document(page_content=text.strip(), metadata={"source": filename, "type": "pptx"}))
                        print(f"Loaded: {filename} ({len(text)} chars)")
                    else:
                        print(f"⚠️ Skipping {filename}: no extractable text")
                # Load Markdown
                elif filename.endswith(".md"):
                    with open(filepath, "r", encoding="utf-8") as f:
                        text = f.read()
                    if text.strip():
                        documents.append(Document(page_content=text.strip(), metadata={"source": filename, "type": "markdown"}))
                        print(f"Loaded: {filename} ({len(text)} chars)")
                    else:
                        print(f"⚠️ Skipping {filename}: empty file")
            except Exception as e: print(f"Error loading {filename}: {e}")
        return documents
    
    def create_vectorstore(self, documents: List[Document]):
        if not documents:
            print("⚠️ No documents to index")
            return
        try:
            # Chunking
            chunks = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).split_documents(documents)
            # Filter out empty chunks that cause embedding errors
            chunks = [c for c in chunks if c.page_content and c.page_content.strip()]
            if not chunks:
                print("⚠️ All chunks were empty after filtering")
                return
            print(f"Created {len(chunks)} chunks (after filtering)")
            # Indexing
            self.vectorstore = FAISS.from_documents(chunks, self.embeddings)
            print("✓ Vector store created successfully")
        except Exception as e:
            print(f"⚠️ Vector store creation failed: {e}")
            import traceback
            traceback.print_exc()
            self.vectorstore = None
        
    def setup_rag_chain(self, language: str = "en"):
        if not self.vectorstore: return
        
        retriever = self.vectorstore.as_retriever(search_kwargs={"k": 5})
        
        # Map language codes to full names for better LLM understanding
        language_map = {
            "en": "English",
            "hi": "Hindi",
            "ta": "Tamil",
            "te": "Telugu"
        }
        language_name = language_map.get(language, "English")
        
        # Multilingual Prompt
        prompt = ChatPromptTemplate.from_template("""You are a helpful AI assistant. Use context to answer.
IMPORTANT:
1. Answer ONLY based on context.
2. If info missing, say: "Information not available in internal documents."
3. Respond in {language}.

Context:
{context}

Question: {question}""")
        
        self.rag_chain = (
            {"context": retriever, "question": RunnablePassthrough(), "language": lambda x: language_name}
            | prompt | self.llm | StrOutputParser()
        )

# Singleton instance
rag_system = None

def initialize_rag():
    global rag_system
    rag_system = LangChainRAG()
    rag_system.create_vectorstore(rag_system.load_documents())

def refresh_documents():
    """Refresh documents from blob storage and rebuild the vector store."""
    global rag_system
    if rag_system:
        print("🔄 Refreshing documents from blob storage...")
        rag_system.create_vectorstore(rag_system.load_documents())
        print("✓ Documents refreshed successfully")

def get_rag_response(query: str, language: str = "en") -> str:
    if not rag_system: initialize_rag()
    
    # Check if RAG is active
    if not rag_system.vectorstore:
        return "System is running in basic mode (no documents indexed). Please add documents to enable RAG."
    
    # Setup chain for current language
    rag_system.setup_rag_chain(language)
    
    try:
        return rag_system.rag_chain.invoke(query)
    except Exception as e:
        print(f"RAG Error: {e}")
        return "Error processing request."
